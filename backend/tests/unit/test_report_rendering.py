"""Report composition and the three renderers.

The assertions worth having here are about *fidelity*: that all three formats
carry the same figures, that a caveat never gets lost on the way into a
document, and that an empty section is rendered with its reason rather than
dropped. A test that only checks bytes came back would pass on a blank page.
"""

from datetime import UTC, date, datetime
from io import BytesIO

import pytest
from openpyxl import load_workbook
from pptx import Presentation

from app.infrastructure.reporting import RENDERERS, renderer_for
from app.services.reporting import commentary
from app.services.reporting.contracts import (
    Block,
    BlockKind,
    ChartKind,
    ExportFormat,
    Kpi,
    Report,
    Section,
)


def _report(sections: tuple[Section, ...] | None = None) -> Report:
    return Report(
        title="Retail Performance Review",
        subtitle="28 days to 21 July 2026",
        period_start=date(2026, 6, 24),
        period_end=date(2026, 7, 21),
        generated_at=datetime(2026, 7, 22, 9, 0, tzinfo=UTC),
        sections=sections if sections is not None else (_kpi_section(), _chart_section()),
        caveats=("Figures come from the surface that owns them.",),
    )


def _kpi_section() -> Section:
    return Section(
        key="kpis",
        title="Key Performance Indicators",
        blocks=(
            Block(
                kind=BlockKind.KPI_GRID,
                kpis=(
                    Kpi("Net Revenue", 10_511_820.0, "currency", comparison=10_770_000.0),
                    Kpi("Margin Rate", 0.412, "rate", comparison=0.405),
                ),
                note="Margin rate is recomputed at the period grain.",
            ),
        ),
    )


def _chart_section() -> Section:
    return Section(
        key="trend",
        title="Performance",
        blocks=(
            Block(
                kind=BlockKind.CHART,
                title="Net Revenue by Region",
                chart_kind=ChartKind.BAR,
                chart_category_column="region",
                chart_value_columns=("net_revenue",),
                columns=("region", "net_revenue"),
                rows=(("Northeast", 1_500_000.0), ("West", 2_100_000.0)),
                note="Bars rather than a line: regions have no inherent order.",
            ),
        ),
    )


# ── Every format renders a real document ─────────────────────────────


@pytest.mark.parametrize("export_format", list(ExportFormat))
def test_every_format_produces_a_valid_document(export_format: ExportFormat) -> None:
    payload = renderer_for(export_format).render(_report())
    assert len(payload) > 1000, "a document this small is a blank page"

    signatures = {
        ExportFormat.PDF: b"%PDF",
        ExportFormat.PPTX: b"PK",  # OOXML is a zip
        ExportFormat.XLSX: b"PK",
    }
    assert payload.startswith(signatures[export_format])


def test_every_format_has_a_renderer() -> None:
    """A format in the enum with no renderer is a 500 waiting to happen."""
    assert set(RENDERERS) == set(ExportFormat)


@pytest.mark.parametrize("export_format", list(ExportFormat))
def test_renderers_declare_a_media_type_and_extension(
    export_format: ExportFormat,
) -> None:
    renderer = renderer_for(export_format)
    assert renderer.media_type
    assert renderer.extension == export_format.value


# ── The formats agree with each other ────────────────────────────────


def test_the_workbook_carries_every_section() -> None:
    """One document, three renderers: none may quietly drop a section."""
    report = _report()
    workbook = load_workbook(BytesIO(RENDERERS[ExportFormat.XLSX].render(report)))

    for section in report.sections:
        assert section.title[:31] in workbook.sheetnames


def test_the_deck_has_a_slide_per_section_plus_covers() -> None:
    report = _report()
    presentation = Presentation(BytesIO(RENDERERS[ExportFormat.PPTX].render(report)))
    # Title slide + one per section + the caveat slide.
    assert len(presentation.slides) == len(report.sections) + 2


def test_charts_are_native_objects_not_pictures() -> None:
    """A picture of a chart cannot be re-sorted, so it gets rebuilt by hand —
    and the hand-built version is the one that disagrees with the report."""
    report = _report()

    workbook = load_workbook(BytesIO(RENDERERS[ExportFormat.XLSX].render(report)))
    assert sum(len(workbook[name]._charts) for name in workbook.sheetnames) >= 1

    presentation = Presentation(BytesIO(RENDERERS[ExportFormat.PPTX].render(report)))
    assert any(shape.has_chart for slide in presentation.slides for shape in slide.shapes)


def test_the_workbook_writes_numbers_as_numbers() -> None:
    """A workbook of numeric-looking text is one nobody can total."""
    workbook = load_workbook(BytesIO(RENDERERS[ExportFormat.XLSX].render(_report())))
    sheet = workbook["Performance"]

    numeric = [
        cell.value
        for row in sheet.iter_rows()
        for cell in row
        if isinstance(cell.value, int | float)
    ]
    assert 1_500_000.0 in numeric


# ── Caveats survive into every format ────────────────────────────────


def _workbook_text(payload: bytes) -> str:
    workbook = load_workbook(BytesIO(payload))
    return " ".join(
        str(cell.value)
        for name in workbook.sheetnames
        for row in workbook[name].iter_rows()
        for cell in row
        if cell.value is not None
    )


def _deck_text(payload: bytes) -> str:
    presentation = Presentation(BytesIO(payload))
    return " ".join(
        shape.text_frame.text
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )


def _caveat_report(marker: str) -> Report:
    return Report(
        title="T",
        subtitle="S",
        period_start=date(2026, 6, 24),
        period_end=date(2026, 7, 21),
        generated_at=datetime(2026, 7, 22, tzinfo=UTC),
        sections=(_kpi_section(),),
        caveats=(marker,) if marker else (),
    )


def test_caveats_reach_the_workbook() -> None:
    """A caveat dropped in rendering is a number acted on unqualified."""
    marker = "UNIQUE-CAVEAT-MARKER"
    assert marker in _workbook_text(RENDERERS[ExportFormat.XLSX].render(_caveat_report(marker)))


def test_caveats_reach_the_deck() -> None:
    marker = "UNIQUE-CAVEAT-MARKER"
    assert marker in _deck_text(RENDERERS[ExportFormat.PPTX].render(_caveat_report(marker)))


def test_caveats_reach_the_pdf() -> None:
    """PDF content streams are compressed, so the text is not greppable.

    Rendering the same report with and without the caveat and comparing sizes
    proves the text was laid out rather than dropped, without adding a PDF
    parser as a test dependency.
    """
    renderer = renderer_for(ExportFormat.PDF)
    with_caveat = renderer.render(_caveat_report("UNIQUE-CAVEAT-MARKER " * 40))
    without = renderer.render(_caveat_report(""))

    assert len(with_caveat) > len(without)


def test_block_notes_reach_the_workbook() -> None:
    workbook = load_workbook(BytesIO(RENDERERS[ExportFormat.XLSX].render(_report())))
    text = " ".join(
        str(cell.value)
        for name in workbook.sheetnames
        for row in workbook[name].iter_rows()
        for cell in row
        if cell.value
    )
    assert "no inherent order" in text


# ── Empty sections are rendered, not dropped ─────────────────────────


def test_an_empty_section_carries_its_reason() -> None:
    """ "Nothing found" and "never ran" are opposite conclusions."""
    reason = "No forecast has been published."
    report = _report(
        sections=(Section(key="forecast", title="Outlook", unavailable_reason=reason),)
    )

    workbook = load_workbook(BytesIO(RENDERERS[ExportFormat.XLSX].render(report)))
    assert "Outlook" in workbook.sheetnames

    presentation = Presentation(BytesIO(RENDERERS[ExportFormat.PPTX].render(report)))
    text = " ".join(
        shape.text_frame.text
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert reason in text


def test_a_report_with_no_sections_still_renders() -> None:
    """A degenerate report must not crash the export."""
    for export_format in ExportFormat:
        assert renderer_for(export_format).render(_report(sections=()))


# ── Renderers survive hostile content ────────────────────────────────


def test_markup_in_data_does_not_break_the_pdf() -> None:
    """reportlab parses inline markup, so an ampersand from a product name
    would otherwise raise mid-render or swallow the rest of the sentence."""
    report = _report(
        sections=(
            Section(
                key="x",
                title="Products",
                blocks=(
                    Block(
                        kind=BlockKind.TABLE,
                        columns=("product",),
                        rows=(("Ben & Jerry's <b>Half Baked</b>",), ("A > B",)),
                    ),
                ),
            ),
        )
    )
    assert renderer_for(ExportFormat.PDF).render(report).startswith(b"%PDF")


def test_illegal_sheet_name_characters_are_stripped() -> None:
    """Excel rejects []:*?/\\ and silently truncates past 31 characters."""
    report = _report(
        sections=(
            Section(
                key="x",
                title="Sales [Q3]: North/South *All* Regions Extended Title",
                blocks=(_kpi_section().blocks[0],),
            ),
        )
    )
    workbook = load_workbook(BytesIO(RENDERERS[ExportFormat.XLSX].render(report)))
    name = workbook.sheetnames[-1]

    assert len(name) <= 31
    assert not set(name) & set(r"[]:*?/\\")


# ── Commentary says only what it can point at ────────────────────────


def test_commentary_names_the_volume_versus_rate_split() -> None:
    """Revenue down while margin rate holds is a volume problem.

    It is invisible in either figure alone, and it changes who owns the fix.
    """
    section = commentary.build(
        kpis=[
            Kpi("Net Revenue", 900.0, "currency", comparison=1000.0),
            Kpi("Margin Rate", 0.42, "rate", comparison=0.40),
        ],
        rca={},
        recommendations={},
        forecast={},
    )
    text = " ".join(block.text for block in section.blocks)
    assert "volume problem" in text


def test_commentary_says_nothing_about_causes_when_none_were_found() -> None:
    """Rather than reaching for a plausible one."""
    section = commentary.build(
        kpis=[Kpi("Net Revenue", 1000.0, "currency", comparison=1000.0)],
        rca={"findings": []},
        recommendations={},
        forecast={},
    )
    text = " ".join(block.text for block in section.blocks)
    assert "driver" not in text.lower()


def test_a_trivial_move_is_described_as_flat() -> None:
    """Calling a 0.4% move a decline trains readers to discount the language."""
    section = commentary.build(
        kpis=[Kpi("Net Revenue", 1004.0, "currency", comparison=1000.0)],
        rca={},
        recommendations={},
        forecast={},
    )
    assert "broadly flat" in section.blocks[0].text


def test_commentary_flags_estimates_resting_on_assumptions() -> None:
    section = commentary.build(
        kpis=[],
        rca={},
        recommendations={
            "recommendations": [
                {
                    "action": "Reorder AC-1010",
                    "owner": "inventory",
                    "impact": {"rests_on_unmeasured_assumptions": True},
                }
            ],
            "gross_profit_opportunity": 100.0,
            "net_profit_opportunity": 100.0,
        },
        forecast={},
    )
    text = " ".join(block.text for block in section.blocks)
    assert "not measured" in text


def test_an_empty_period_produces_honest_commentary() -> None:
    section = commentary.build(kpis=[], rca={}, recommendations={}, forecast={})
    assert "little to say" in section.blocks[0].text


def test_the_executive_summary_reports_when_nothing_cleared_a_threshold() -> None:
    section = commentary.executive_summary(
        kpis=[], rca={}, recommendations={}, forecast={}, period_label="July"
    )
    bullets = [bullet for block in section.blocks for bullet in block.bullets]
    assert any("finding, not a gap" in bullet for bullet in bullets)


# ── KPI arithmetic ───────────────────────────────────────────────────


def test_kpi_change_is_relative_to_the_comparison() -> None:
    assert Kpi("R", 110.0, "currency", comparison=100.0).change == pytest.approx(0.10)


def test_kpi_change_is_absent_rather_than_infinite_without_a_baseline() -> None:
    assert Kpi("R", 110.0, "currency").change is None
    assert Kpi("R", 110.0, "currency", comparison=0.0).change is None


def test_rates_and_currency_format_differently() -> None:
    assert Kpi("Margin Rate", 0.412, "rate").formatted() == "41.2%"
    assert Kpi("Revenue", 1_234_567.0, "currency").formatted() == "1,234,567"
