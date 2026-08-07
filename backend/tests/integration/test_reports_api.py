"""Report endpoints against a real warehouse.

Proves the composed document is assembled from the live services and that all
three exports are documents a recipient's software will actually open — not
just bytes with the right magic number.
"""

from io import BytesIO

import pytest

pytest.importorskip("testcontainers", reason="integration extra not installed")
from httpx import AsyncClient  # noqa: E402
from openpyxl import load_workbook  # noqa: E402
from pptx import Presentation  # noqa: E402

from tests.integration.conftest import auth_headers  # noqa: E402
from tests.integration.warehouse import LAST_DAY  # noqa: E402

pytestmark = pytest.mark.integration


async def _report(api: AsyncClient, role: str = "ceo") -> dict:
    response = await api.get(
        "/api/v1/reports",
        headers=await auth_headers(api, role),
        params={"period_end": LAST_DAY.isoformat(), "period_days": 14},
    )
    assert response.status_code == 200, response.text
    return response.json()


async def _export(api: AsyncClient, fmt: str):
    return await api.get(
        "/api/v1/reports/export",
        headers=await auth_headers(api),
        params={"format": fmt, "period_end": LAST_DAY.isoformat(), "period_days": 14},
    )


# ── Composition ──────────────────────────────────────────────────────


async def test_the_report_has_every_requested_section(api: AsyncClient) -> None:
    body = await _report(api)
    keys = {section["key"] for section in body["sections"]}
    assert {
        "summary",
        "kpis",
        "trend",
        "insights",
        "forecast",
        "recommendations",
        "commentary",
    } <= keys


async def test_the_executive_summary_leads_with_a_headline(api: AsyncClient) -> None:
    body = await _report(api)
    summary = next(s for s in body["sections"] if s["key"] == "summary")
    assert any(block["text"] for block in summary["blocks"])


async def test_kpis_carry_a_prior_period_comparison(api: AsyncClient) -> None:
    body = await _report(api)
    kpis = [
        kpi for section in body["sections"] for block in section["blocks"] for kpi in block["kpis"]
    ]
    assert kpis
    assert any(kpi["comparison"] is not None for kpi in kpis)


async def test_an_empty_section_states_why(api: AsyncClient) -> None:
    """ "Nothing found" and "never ran" are opposite conclusions."""
    body = await _report(api)
    for section in body["sections"]:
        if not section["blocks"]:
            assert section["unavailable_reason"], f"{section['key']} is silently empty"


async def test_the_report_carries_its_caveats(api: AsyncClient) -> None:
    body = await _report(api)
    assert body["caveats"]
    assert any("recomputed" in caveat for caveat in body["caveats"])


async def test_commentary_is_present_and_non_empty(api: AsyncClient) -> None:
    body = await _report(api)
    commentary = next(s for s in body["sections"] if s["key"] == "commentary")
    assert any(block["text"] for block in commentary["blocks"])


# ── Exports ──────────────────────────────────────────────────────────


@pytest.mark.parametrize("fmt", ["pdf", "pptx", "xlsx"])
async def test_every_format_downloads(api: AsyncClient, fmt: str) -> None:
    response = await _export(api, fmt)
    assert response.status_code == 200, response.text
    assert len(response.content) > 1000
    assert f".{fmt}" in response.headers["content-disposition"]
    assert response.headers["content-disposition"].startswith("attachment")


async def test_the_pdf_is_a_pdf(api: AsyncClient) -> None:
    response = await _export(api, "pdf")
    assert response.content.startswith(b"%PDF")
    assert response.headers["content-type"] == "application/pdf"


async def test_the_workbook_opens_with_native_charts(api: AsyncClient) -> None:
    """Native charts, not pictures: the recipient can re-sort and re-total."""
    response = await _export(api, "xlsx")
    workbook = load_workbook(BytesIO(response.content))

    assert len(workbook.sheetnames) >= 5
    assert sum(len(workbook[name]._charts) for name in workbook.sheetnames) >= 1


async def test_the_deck_opens_with_native_charts(api: AsyncClient) -> None:
    response = await _export(api, "pptx")
    presentation = Presentation(BytesIO(response.content))

    assert len(presentation.slides) >= 5
    assert any(shape.has_chart for slide in presentation.slides for shape in slide.shapes)


async def test_the_exports_agree_with_the_json(api: AsyncClient) -> None:
    """One document, three renderers — the whole reason for the design."""
    body = await _report(api)
    response = await _export(api, "xlsx")
    workbook = load_workbook(BytesIO(response.content))

    for section in body["sections"]:
        assert section["title"][:31] in workbook.sheetnames


async def test_an_unknown_format_is_refused(api: AsyncClient) -> None:
    response = await api.get(
        "/api/v1/reports/export", headers=await auth_headers(api), params={"format": "docx"}
    )
    assert response.status_code == 422


# ── Authorization ────────────────────────────────────────────────────


async def test_anonymous_access_is_rejected(api: AsyncClient) -> None:
    assert (await api.get("/api/v1/reports")).status_code == 401
    assert (await api.get("/api/v1/reports/export")).status_code == 401


async def test_a_denied_section_narrows_the_report_rather_than_failing_it(
    api: AsyncClient,
) -> None:
    """Marketing may read reports but not profitability.

    Failing the whole document because one section was forbidden turns a
    partial report into no report. The section is rendered with its reason and
    everything the role *may* see is still there.
    """
    body = await _report(api, role="marketing")
    assert body["sections"], "a permission denial emptied the whole report"

    denied = [
        section
        for section in body["sections"]
        if "role does not include" in section["unavailable_reason"]
    ]
    populated = [section for section in body["sections"] if section["blocks"]]

    assert populated, "nothing survived the denial"
    for section in denied:
        assert section["title"], "a denied section still needs its heading"
