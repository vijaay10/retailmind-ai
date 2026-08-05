"""PowerPoint export.

One section per slide, and charts as **native PowerPoint chart objects**
rather than images. That matters more here than anywhere else: a deck gets
forwarded, edited, and re-presented, and the first thing someone does is
restyle a chart to match their template or extend a series. A picture cannot
be edited, so it gets deleted and rebuilt by hand from the underlying numbers
— which is where the version that disagrees with this report comes from.

Content is capped hard. A slide holding a twelve-row table is a spreadsheet
someone has pasted into a presentation, and the audience reads none of it.
"""

from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.presentation import Presentation as PresentationType
from pptx.util import Inches, Pt

from app.services.reporting.contracts import (
    Block,
    BlockKind,
    ChartKind,
    ExportFormat,
    Renderer,
    Report,
    Section,
)

# python-pptx ships type hints but leaves several constructors unannotated.
# Ignored at the call site rather than by relaxing the module, so everything
# else here stays under strict checking.
BRAND = RGBColor(0x1F, 0x38, 0x64)  # type: ignore[no-untyped-call]
MUTED = RGBColor(0x59, 0x59, 0x59)  # type: ignore[no-untyped-call]

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

#: What actually fits on a slide someone reads from the back of a room.
MAX_TABLE_ROWS = 6
MAX_TABLE_COLUMNS = 5
MAX_BULLETS = 5
MAX_CHART_POINTS = 30

_CHART_TYPES = {
    ChartKind.LINE: XL_CHART_TYPE.LINE_MARKERS,
    ChartKind.BAR: XL_CHART_TYPE.COLUMN_CLUSTERED,
    ChartKind.HORIZONTAL_BAR: XL_CHART_TYPE.BAR_CLUSTERED,
}


class PptxRenderer(Renderer):
    format = ExportFormat.PPTX
    media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    extension = "pptx"

    def render(self, report: Report) -> bytes:
        presentation = Presentation()
        presentation.slide_width = SLIDE_WIDTH
        presentation.slide_height = SLIDE_HEIGHT

        self._title_slide(presentation, report)
        for section in report.sections:
            self._section_slide(presentation, section)
        self._caveat_slide(presentation, report)

        buffer = BytesIO()
        presentation.save(buffer)
        return buffer.getvalue()

    # ── Slides ───────────────────────────────────────────────────────

    def _title_slide(self, presentation: PresentationType, report: Report) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._text(
            slide,
            report.title,
            Inches(0.9),
            Inches(2.4),
            Inches(11.5),
            Inches(1.2),
            size=40,
            bold=True,
            colour=BRAND,
        )
        self._text(
            slide,
            report.subtitle,
            Inches(0.9),
            Inches(3.6),
            Inches(11.5),
            Inches(0.8),
            size=16,
            colour=MUTED,
        )
        self._text(
            slide,
            f"{report.period_label}  ·  generated {report.generated_at:%d %b %Y %H:%M} UTC",
            Inches(0.9),
            Inches(4.4),
            Inches(11.5),
            Inches(0.5),
            size=12,
            colour=MUTED,
        )

    def _section_slide(self, presentation: PresentationType, section: Section) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._text(
            slide,
            section.title,
            Inches(0.6),
            Inches(0.4),
            Inches(12),
            Inches(0.8),
            size=28,
            bold=True,
            colour=BRAND,
        )
        if section.subtitle:
            self._text(
                slide,
                section.subtitle,
                Inches(0.6),
                Inches(1.15),
                Inches(12),
                Inches(0.5),
                size=13,
                colour=MUTED,
            )

        if section.unavailable_reason:
            self._text(
                slide,
                section.unavailable_reason,
                Inches(0.6),
                Inches(2.2),
                Inches(11.5),
                Inches(2),
                size=14,
                colour=MUTED,
            )
            return

        top = Inches(1.8)
        for block in section.blocks:
            top = self._block(slide, block, top)

    def _block(self, slide: Any, block: Block, top: Any) -> Any:
        if top > Inches(6.4):
            return top  # the slide is full; the detail lives in the workbook

        if block.kind is BlockKind.NARRATIVE and block.text:
            self._text(slide, block.text, Inches(0.6), top, Inches(12), Inches(1.4), size=14)
            return top + Inches(1.5)

        if block.kind is BlockKind.CALLOUT:
            for bullet in block.bullets[:MAX_BULLETS]:
                self._text(
                    slide, f"•  {bullet}", Inches(0.8), top, Inches(11.6), Inches(0.55), size=12
                )
                top += Inches(0.55)
            return top + Inches(0.2)

        if block.kind is BlockKind.KPI_GRID and block.kpis:
            return self._kpis(slide, block, top)

        if block.chart_kind and block.rows:
            return self._chart(slide, block, top)

        if block.rows:
            return self._table(slide, block, top)

        return top

    def _kpis(self, slide: Any, block: Block, top: Any) -> Any:
        """KPI tiles across the slide."""
        tiles = block.kpis[:5]
        width = Inches(2.3)
        for index, kpi in enumerate(tiles):
            left = Inches(0.6) + index * Inches(2.45)
            self._text(slide, kpi.label, left, top, width, Inches(0.4), size=11, colour=MUTED)
            self._text(
                slide,
                kpi.formatted(),
                left,
                top + Inches(0.35),
                width,
                Inches(0.7),
                size=24,
                bold=True,
                colour=BRAND,
            )
            if kpi.change is not None:
                self._text(
                    slide,
                    f"{kpi.change:+.1%} {kpi.comparison_label}",
                    left,
                    top + Inches(1.0),
                    width,
                    Inches(0.4),
                    size=10,
                    colour=MUTED,
                )
        return top + Inches(1.6)

    def _chart(self, slide: Any, block: Block, top: Any) -> Any:
        """A real chart object, not a picture of one."""
        rows = block.rows[:MAX_CHART_POINTS]
        categories = [str(row[0]) for row in rows]

        data = CategoryChartData()  # type: ignore[no-untyped-call]
        data.categories = categories
        for offset, name in enumerate(block.chart_value_columns, start=1):
            if offset >= len(block.columns):
                continue
            data.add_series(  # type: ignore[no-untyped-call]
                name.replace("_", " ").title(),
                tuple(_number(row[offset]) if offset < len(row) else 0.0 for row in rows),
            )

        chart_type = (
            _CHART_TYPES.get(block.chart_kind, XL_CHART_TYPE.COLUMN_CLUSTERED)
            if block.chart_kind is not None
            else XL_CHART_TYPE.COLUMN_CLUSTERED
        )
        graphic = slide.shapes.add_chart(
            chart_type,
            Inches(0.6),
            top,
            Inches(12),
            Inches(4.2),
            data,
        )
        chart = graphic.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        return top + Inches(4.4)

    def _table(self, slide: Any, block: Block, top: Any) -> Any:
        columns = list(block.columns[:MAX_TABLE_COLUMNS])
        rows = block.rows[:MAX_TABLE_ROWS]

        shape = slide.shapes.add_table(
            len(rows) + 1,
            len(columns),
            Inches(0.6),
            top,
            Inches(12),
            Inches(0.4) * (len(rows) + 1),
        )
        table = shape.table

        for index, name in enumerate(columns):
            cell = table.cell(0, index)
            cell.text = name.replace("_", " ").title()
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].runs[0].font.bold = True

        for row_index, record in enumerate(rows, start=1):
            for column_index in range(len(columns)):
                value = record[column_index] if column_index < len(record) else ""
                cell = table.cell(row_index, column_index)
                cell.text = _cell(value)
                cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)

        return top + Inches(0.4) * (len(rows) + 1) + Inches(0.3)

    def _caveat_slide(self, presentation: PresentationType, report: Report) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._text(
            slide,
            "How to read this report",
            Inches(0.6),
            Inches(0.5),
            Inches(12),
            Inches(0.8),
            size=26,
            bold=True,
            colour=BRAND,
        )
        top: int = Inches(1.7)
        for caveat in report.caveats:
            self._text(slide, f"•  {caveat}", Inches(0.8), top, Inches(11.6), Inches(1.0), size=13)
            top += Inches(1.0)

    # ── Helpers ──────────────────────────────────────────────────────

    def _text(
        self,
        slide: Any,
        text: str,
        left: Any,
        top: Any,
        width: Any,
        height: Any,
        *,
        size: int = 14,
        bold: bool = False,
        colour: RGBColor | None = None,
    ) -> None:
        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.word_wrap = True
        paragraph = frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        if colour is not None:
            run.font.color.rgb = colour


def _number(value: Any) -> float:
    if not isinstance(value, int | float | str):
        return 0.0
    try:
        return float(value)
    except (TypeError, ValueError):
        return 0.0


def _cell(value: Any) -> str:
    if isinstance(value, float):
        return f"{value:,.0f}" if abs(value) >= 1000 else f"{value:,.2f}"
    if isinstance(value, int):
        return f"{value:,}"
    return str(value)[:48]
